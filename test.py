# Requirements:
# streamlit
# pandas
# matplotlib
# seaborn
# openpyxl
# python-pptx
# pillow
# comtypes (for pptx preview on Windows, optional)
# numpy

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import numpy as np
from PIL import Image
import tempfile

# === Placeholders for advanced modules ===
# (Implement each in its own function for modularity)

def data_upload_module():
    st.header("Step 1: Upload Data")
    uploaded_files = st.file_uploader(
        "Upload Excel (.xlsx) or CSV (.csv) files", type=["xlsx", "csv"], accept_multiple_files=True
    )
    return uploaded_files

def data_cleaning_module(df):
    st.header("Step 2: Data Cleaning Suggestions")
    return df

def data_analysis_module(df):
    st.header("Step 3: Data Analysis & Insights")
    summary = {}
    insights = []
    num_cols = df.select_dtypes(include='number').columns
    cat_cols = df.select_dtypes(include='object').columns

    # Numerical summary
    if len(num_cols) > 0:
        desc = df[num_cols].describe().T
        summary['numerical'] = desc
        for col in num_cols:
            col_data = df[col].dropna()
            if len(col_data) == 0:
                continue
            mean = col_data.mean()
            median = col_data.median()
            std = col_data.std()
            min_ = col_data.min()
            max_ = col_data.max()
            outliers = col_data[(col_data < mean - 3*std) | (col_data > mean + 3*std)]
            if len(outliers) > 0:
                insights.append(f"Column '{col}' has {len(outliers)} outlier(s).")
            insights.append(f"Column '{col}': mean={mean:.2f}, median={median:.2f}, min={min_}, max={max_}.")
    # Categorical summary
    if len(cat_cols) > 0:
        cat_summary = {}
        for col in cat_cols:
            vc = df[col].value_counts()
            cat_summary[col] = vc
            top = vc.index[0] if not vc.empty else None
            if top:
                insights.append(f"Most frequent value in '{col}' is '{top}' ({vc.iloc[0]} times).")
        summary['categorical'] = cat_summary
    return summary, insights

def plot_and_save(df, sheet_name, chart_selections=None):
    images = []
    captions = []
    num_cols = df.select_dtypes(include='number').columns
    cat_cols = df.select_dtypes(include='object').columns

    # Use chart_selections if provided, else default to all columns
    selected_num_cols = chart_selections.get("num_cols", num_cols) if chart_selections else num_cols
    selected_cat_cols = chart_selections.get("cat_cols", cat_cols) if chart_selections else cat_cols
    chart_types = chart_selections.get("chart_types", ["hist", "line", "bar", "pie"]) if chart_selections else ["hist", "line", "bar", "pie"]

    # Numerical columns: histograms and line plots
    for col in selected_num_cols:
        if "hist" in chart_types:
            fig, ax = plt.subplots()
            sns.histplot(df[col].dropna(), kde=True, ax=ax)
            ax.set_title(f"Distribution of {col}")
            buf = BytesIO()
            plt.savefig(buf, format='png')
            buf.seek(0)
            images.append(buf.read())
            captions.append(f"Distribution of {col} in {sheet_name}.")
            plt.close(fig)
        if "line" in chart_types and pd.api.types.is_datetime64_any_dtype(df.index):
            fig, ax = plt.subplots()
            df[col].plot(ax=ax)
            ax.set_title(f"{col} over time")
            buf = BytesIO()
            plt.savefig(buf, format='png')
            buf.seek(0)
            images.append(buf.read())
            captions.append(f"Trend of {col} over time in {sheet_name}.")
            plt.close(fig)

    # Categorical columns: bar plots and pie charts
    for col in selected_cat_cols:
        vc = df[col].value_counts().head(10)
        if len(vc) > 1:
            if "bar" in chart_types:
                fig, ax = plt.subplots()
                sns.barplot(x=vc.values, y=vc.index, ax=ax)
                ax.set_title(f"Top categories in {col}")
                buf = BytesIO()
                plt.savefig(buf, format='png')
                buf.seek(0)
                images.append(buf.read())
                captions.append(f"Top categories in {col} in {sheet_name}.")
                plt.close(fig)
            if "pie" in chart_types:
                fig, ax = plt.subplots()
                vc.plot.pie(autopct='%1.1f%%', ax=ax)
                ax.set_ylabel('')
                ax.set_title(f"Category distribution in {col}")
                buf = BytesIO()
                plt.savefig(buf, format='png')
                buf.seek(0)
                images.append(buf.read())
                captions.append(f"Category distribution in {col} in {sheet_name}.")
                plt.close(fig)
    return images, captions

def add_title_slide(prs, file_name):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Data Analysis Report"
    subtitle.text = f"File: {file_name}\nDate: {datetime.now().strftime('%Y-%m-%d')}"

def add_chart_slide(prs, image_bytes, caption):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = caption
    left = Inches(1)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=Inches(6))
    # Optionally add caption as text box
    # ...

def add_summary_slide(prs, insights):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Findings"
    tf = slide.placeholders[1].text_frame
    for insight in insights:
        p = tf.add_paragraph()
        p.text = insight
        p.level = 0

def create_pptx(file_name, all_images, all_captions, insights):
    prs = Presentation()
    add_title_slide(prs, file_name)
    for img, cap in zip(all_images, all_captions):
        add_chart_slide(prs, img, cap)
    add_summary_slide(prs, insights)
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def reporting_module(selected_charts, insights):
    st.header("Step 5: Reporting & Export")
    st.info("Reporting module is under development.")
    # Placeholder for reporting functionality

def collaboration_module():
    st.header("Step 6: Collaboration & Sharing")
    st.info("Collaboration module is under development.")
    # Placeholder for collaboration functionality

def ai_ml_module(df):
    st.header("Step 7: AI/ML Features")
    st.info("AI/ML module is under development.")
    # Placeholder for AI/ML functionality

def analyze_and_clean(df, inplace=False, fill_missing='auto', drop_threshold=0.5, return_flagged=True):
    """
    Perform ground-level data analysis and cleanup on a pandas DataFrame.
    Args:
        df: pandas DataFrame
        inplace: if True, modify df in place; else work on a copy
        fill_missing: 'auto', 'mean', 'median', 'mode', 'drop', or None
        drop_threshold: if fraction of missing in a column > threshold, drop column
        return_flagged: if True, return (cleaned_df, flagged_df)
    Returns:
        cleaned_df, flagged_df (if return_flagged), else just cleaned_df
    """
    if not inplace:
        df = df.copy()
    print("=== Initial DataFrame shape:", df.shape)

    flagged_rows = pd.DataFrame(index=df.index)
    # 1. Check for missing values per column
    print("\n=== Missing Values Per Column ===")
    missing = df.isnull().sum()
    print(missing)
    # Suggest drop columns with too many missing
    cols_to_drop = missing[missing / len(df) > drop_threshold].index.tolist()
    if cols_to_drop:
        print(f"\nColumns with >{int(drop_threshold*100)}% missing values: {cols_to_drop} (will be dropped)")
        df.drop(columns=cols_to_drop, inplace=True)
    # Fill or drop missing values
    for col in df.columns:
        if df[col].isnull().any():
            if fill_missing == 'auto':
                if df[col].dtype.kind in 'biufc':  # numeric
                    fill_val = df[col].median()
                    print(f"Filling missing in '{col}' with median: {fill_val}")
                    df[col].fillna(fill_val, inplace=True)
                elif np.issubdtype(df[col].dtype, np.datetime64):
                    fill_val = df[col].mode().iloc[0] if not df[col].mode().empty else None
                    print(f"Filling missing in '{col}' with mode: {fill_val}")
                    df[col].fillna(fill_val, inplace=True)
                else:
                    fill_val = df[col].mode().iloc[0] if not df[col].mode().empty else ""
                    print(f"Filling missing in '{col}' with mode: {fill_val}")
                    df[col].fillna(fill_val, inplace=True)
            elif fill_missing == 'drop':
                print(f"Dropping rows with missing in '{col}'")
                df = df[df[col].notnull()]
            elif fill_missing in ['mean', 'median', 'mode']:
                if fill_missing == 'mean' and df[col].dtype.kind in 'biufc':
                    fill_val = df[col].mean()
                elif fill_missing == 'median' and df[col].dtype.kind in 'biufc':
                    fill_val = df[col].median()
                else:
                    fill_val = df[col].mode().iloc[0] if not df[col].mode().empty else ""
                print(f"Filling missing in '{col}' with {fill_missing}: {fill_val}")
                df[col].fillna(fill_val, inplace=True)
    # Flag rows with any missing values (after filling)
    flagged_rows['missing'] = df.isnull().any(axis=1)

    # 2. Remove duplicate rows
    before = len(df)
    df_nodup = df.drop_duplicates()
    after = len(df_nodup)
    print(f"\n=== Duplicate Removal ===\nRemoved {before - after} duplicate rows.")
    flagged_rows['duplicate'] = ~df.index.isin(df_nodup.index)
    df = df_nodup

    # 3. Convert column data types
    print("\n=== Data Type Conversion ===")
    for col in df.columns:
        orig_dtype = df[col].dtype
        # Try numeric
        if df[col].dtype == object:
            try:
                df[col] = pd.to_numeric(df[col])
                print(f"Converted '{col}' to numeric.")
            except Exception:
                # Try datetime
                try:
                    df[col] = pd.to_datetime(df[col])
                    print(f"Converted '{col}' to datetime.")
                except Exception:
                    pass
        # Try datetime for columns with 'date' in name
        if 'date' in col.lower() and not np.issubdtype(df[col].dtype, np.datetime64):
            try:
                df[col] = pd.to_datetime(df[col])
                print(f"Converted '{col}' to datetime (by name).")
            except Exception:
                pass
        if orig_dtype != df[col].dtype:
            print(f"Column '{col}': {orig_dtype} -> {df[col].dtype}")

    # 4. Standardize string columns
    print("\n=== Standardizing String Columns ===")
    for col in df.select_dtypes(include='object').columns:
        df[col] = df[col].astype(str).str.strip().str.lower()
        print(f"Standardized '{col}' (strip, lower).")

    # 5. Print summaries
    print("\n=== Column Data Types ===")
    print(df.dtypes)
    print("\n=== Basic Statistics ===")
    print(df.describe(include='all', datetime_is_numeric=True))
    print("\n=== Value Counts for Categorical Columns ===")
    for col in df.select_dtypes(include='object').columns:
        print(f"\nValue counts for '{col}':")
        print(df[col].value_counts())

    # 6. Detect and highlight outliers (IQR method)
    print("\n=== Outlier Detection (IQR method) ===")
    outlier_flags = pd.DataFrame(False, index=df.index, columns=df.select_dtypes(include=np.number).columns)
    for col in df.select_dtypes(include=np.number).columns:
        Q1 = df[col].quantile(0.25)
        Q3 = df[col].quantile(0.75)
        IQR = Q3 - Q1
        lower = Q1 - 1.5 * IQR
        upper = Q3 + 1.5 * IQR
        outliers = (df[col] < lower) | (df[col] > upper)
        n_out = outliers.sum()
        print(f"Column '{col}': {n_out} outlier(s) flagged.")
        outlier_flags[col] = outliers
    flagged_rows['outlier'] = outlier_flags.any(axis=1)

    # 7. Flag suspicious rows
    flagged_rows['suspicious'] = flagged_rows.any(axis=1)
    print(f"\n=== Suspicious Rows: {flagged_rows['suspicious'].sum()} flagged ===")

    # 8. Print before/after row counts
    print(f"\n=== Row Count: Before: {before}, After: {len(df)} ===")

    # 9. Return cleaned DataFrame and flagged DataFrame
    if return_flagged:
        flagged_df = df[flagged_rows['suspicious']]
        return df, flagged_df
    else:
        return df

def generate_clear_chart(
    df,
    chart_type="bar",
    x=None,
    y=None,
    title=None,
    save_path=None,
    return_fig=False,
    max_pie_slices=8
):
    """
    Generate a clear, simple, and sorted chart from a pandas DataFrame.
    Supports: bar, line, pie.
    """
    import matplotlib.pyplot as plt
    import seaborn as sns
    from matplotlib.ticker import FuncFormatter
    import matplotlib.dates as mdates

    # --- Style ---
    plt.rcParams.update({
        "axes.titlesize": 18,
        "axes.labelsize": 14,
        "xtick.labelsize": 13,
        "ytick.labelsize": 13,
        "axes.titlepad": 18,
        "axes.labelpad": 10,
        "axes.edgecolor": "#aaa",
        "axes.linewidth": 1.2,
        "figure.facecolor": "#fff",
        "axes.facecolor": "#fff",
        "legend.fontsize": 13,
        "font.family": "DejaVu Sans"
    })
    fig, ax = plt.subplots(figsize=(8, 5))

    def auto_title(chart_type, x, y):
        if chart_type == "bar":
            return f"{y or 'Value'} by {x or 'Category'}"
        elif chart_type == "line":
            return f"{y or 'Value'} over {x or 'Time'}"
        elif chart_type == "pie":
            return f"{x or 'Category'} Distribution"
        else:
            return "Chart"

    # --- Bar Chart ---
    if chart_type == "bar":
        if x is None or y is None:
            raise ValueError("x and y must be specified for bar chart")
        # Group and sort
        data = df.groupby(x, dropna=False)[y].sum().sort_values(ascending=False)
        data = data.head(15)  # Top 15 for clarity
        bars = ax.barh(
            data.index.astype(str),
            data.values,
            color=sns.color_palette("Blues_r", len(data)),
            edgecolor="#333",
            height=0.6
        )
        ax.set_xlabel(y)
        ax.set_ylabel(x)
        ax.set_title(title or auto_title("bar", x, y))
        ax.invert_yaxis()
        for bar in bars:
            width = bar.get_width()
            ax.text(width + max(data.values)*0.01, bar.get_y() + bar.get_height()/2,
                    f"{width:,.0f}", va="center", fontsize=13, color="#222")
        ax.spines[['top', 'right', 'left']].set_visible(False)
        ax.xaxis.set_major_formatter(FuncFormatter(lambda x, _: f"{int(x):,}"))
        ax.grid(axis='x', linestyle='--', alpha=0.3)
        plt.tight_layout()

    # --- Line Chart ---
    elif chart_type == "line":
        if x is None or y is None:
            raise ValueError("x and y must be specified for line chart")
        x_data = pd.to_datetime(df[x], errors="coerce")
        sorted_df = df.assign(_x=x_data).sort_values("_x")
        ax.plot(sorted_df["_x"], sorted_df[y], marker="o", color="#4F8DFD", linewidth=2)
        ax.set_xlabel(x)
        ax.set_ylabel(y)
        ax.set_title(title or auto_title("line", x, y))
        ax.spines[['top', 'right']].set_visible(False)
        ax.grid(axis='y', linestyle='--', alpha=0.3)
        # Date formatting
        if pd.api.types.is_datetime64_any_dtype(sorted_df["_x"]):
            locator = mdates.AutoDateLocator()
            formatter = mdates.ConciseDateFormatter(locator)
            ax.xaxis.set_major_locator(locator)
            ax.xaxis.set_major_formatter(formatter)
        plt.setp(ax.xaxis.get_majorticklabels(), rotation=30, ha="right")
        plt.tight_layout()

    # --- Pie Chart ---
    elif chart_type == "pie":
        if x is None:
            raise ValueError("x must be specified for pie chart")
        counts = df[x].value_counts(dropna=False)
        if max_pie_slices and len(counts) > max_pie_slices:
            top = counts[:max_pie_slices-1]
            others = counts[max_pie_slices-1:].sum()
            counts = top.append(pd.Series({"Others": others}))
        colors = sns.color_palette("pastel", len(counts))
        wedges, texts, autotexts = ax.pie(
            counts.values,
            labels=[str(i) for i in counts.index],
            autopct=lambda pct: f"{pct:.1f}%",
            startangle=90,
            colors=colors,
            textprops={'fontsize': 13, 'color': "#222"}
        )
        ax.set_title(title or auto_title("pie", x, None))
        plt.tight_layout()
    else:
        raise ValueError("chart_type must be one of: 'bar', 'line', 'pie'")

    # --- Clean up ---
    ax.set_facecolor("#fff")
    fig.patch.set_facecolor("#fff")
    plt.subplots_adjust(left=0.18, right=0.98, top=0.88, bottom=0.18)

    # Save or return
    if save_path:
        plt.savefig(save_path, bbox_inches="tight", dpi=150)
    if return_fig:
        return fig, ax
    else:
        plt.show()
        plt.close(fig)

def pptx_to_slide_images(pptx_bytes):
    """Convert pptx bytes to a list of slide images (PNG)."""
    import comtypes.client
    import os
    import uuid
    import shutil

    # Only works on Windows with PowerPoint installed
    # Save pptx to temp file
    tmp_dir = tempfile.mkdtemp()
    pptx_path = os.path.join(tmp_dir, f"{uuid.uuid4()}.pptx")
    with open(pptx_path, "wb") as f:
        f.write(pptx_bytes)
    # Export slides as PNG
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    presentation = powerpoint.Presentations.Open(pptx_path)
    export_dir = os.path.join(tmp_dir, "slides")
    os.makedirs(export_dir, exist_ok=True)
    presentation.SaveAs(export_dir, 17)  # 17 = ppSaveAsPNG
    presentation.Close()
    powerpoint.Quit()
    # Collect slide images
    slide_imgs = []
    for fname in sorted(os.listdir(export_dir)):
        if fname.endswith(".PNG"):
            with open(os.path.join(export_dir, fname), "rb") as imgf:
                slide_imgs.append(imgf.read())
    shutil.rmtree(tmp_dir)
    return slide_imgs

def main():
    st.set_page_config(page_title="Data-Driven SaaS Assistant", layout="wide")
    st.title("Excel Data-Driven SaaS Assistant")

    # Sidebar for upload and global options
    with st.sidebar:
        st.header("Upload & Settings")
        uploaded_file = st.file_uploader("Upload Excel (.xlsx) or CSV (.csv)", type=["xlsx", "csv"])
        st.markdown("---")
        st.info("Configure your analysis and report below.")

    if uploaded_file is not None:
        try:
            file_name = uploaded_file.name
            file_bytes = uploaded_file.read()
            if not file_bytes:
                st.error("Uploaded file is empty.")
                return
            is_excel = file_name.lower().endswith(".xlsx")
            is_csv = file_name.lower().endswith(".csv")
            if is_excel or is_csv:
                # Use columns for main workflow
                col1, col2 = st.columns([2, 1])
                with col1:
                    st.subheader("Data Preview & Analysis")
                with col2:
                    st.subheader("Configuration")

                # Data processing and configuration in sidebar
                with st.sidebar:
                    st.markdown("### Report Sections")
                    show_data = st.checkbox("Show Data Table", value=True)
                    show_summary = st.checkbox("Show Summary Statistics", value=True)
                    show_charts = st.checkbox("Show Charts", value=True)
                    show_qa = st.checkbox("Enable Data Q&A", value=True)
                    show_report_edit = st.checkbox("Enable Report Editing", value=True)

                # Main content in columns
                with col1:
                    all_images, all_captions, all_insights, summary, insights, images, captions = process_uploaded_data(
                        file_name, file_bytes, is_excel,
                        show_data=show_data,
                        show_summary=show_summary,
                        show_charts=show_charts,
                        show_qa=show_qa
                    )

                # Report editing and generation in right column
                with col2:
                    if show_report_edit:
                        st.subheader("Edit Report Structure")
                        # --- Editable Report Structure Before Generation ---
                        # Prepare editable structure: slides = [{type, title, content/image/caption}]
                        slides = []

                        # Title slide
                        title_slide = {
                            "type": "title",
                            "title": "Data Analysis Report",
                            "subtitle": f"File: {file_name}\nDate: {datetime.now().strftime('%Y-%m-%d')}"
                        }
                        slides.append(title_slide)

                        # Chart slides
                        for img, cap in zip(all_images, all_captions):
                            slides.append({
                                "type": "chart",
                                "caption": cap,
                                "image": img
                            })

                        # Key findings slide
                        key_findings = all_insights if is_excel else (insights if insights else [])
                        slides.append({
                            "type": "summary",
                            "title": "Key Findings",
                            "insights": key_findings
                        })

                        # Editable UI for slides
                        new_slides = []
                        st.markdown("**You can edit, remove, or reorder slides below.**")
                        for idx, slide in enumerate(slides):
                            with st.expander(f"Slide {idx+1}: {slide['type'].capitalize()}"):
                                remove = st.checkbox("Remove this slide", key=f"remove_slide_{idx}")
                                if remove:
                                    continue
                                if slide["type"] == "title":
                                    title = st.text_input("Title", value=slide["title"], key=f"title_{idx}")
                                    subtitle = st.text_area("Subtitle", value=slide["subtitle"], key=f"subtitle_{idx}")
                                    new_slides.append({
                                        "type": "title",
                                        "title": title,
                                        "subtitle": subtitle
                                    })
                                elif slide["type"] == "chart":
                                    caption = st.text_area("Caption", value=slide["caption"], key=f"caption_{idx}")
                                    st.image(slide["image"], caption="Chart Preview")
                                    new_slides.append({
                                        "type": "chart",
                                        "caption": caption,
                                        "image": slide["image"]
                                    })
                                elif slide["type"] == "summary":
                                    title = st.text_input("Summary Slide Title", value=slide["title"], key=f"summary_title_{idx}")
                                    insights_text = st.text_area(
                                        "Key Findings (one per line)",
                                        value="\n".join(slide["insights"]),
                                        height=200,
                                        key=f"insights_{idx}"
                                    )
                                    insights_list = [line.strip() for line in insights_text.split("\n") if line.strip()]
                                    new_slides.append({
                                        "type": "summary",
                                        "title": title,
                                        "insights": insights_list
                                    })

                        if st.button("Generate PowerPoint Report"):
                            # Build PPTX from edited structure
                            prs = Presentation()
                            # Title slide
                            for slide in new_slides:
                                if slide["type"] == "title":
                                    slide_layout = prs.slide_layouts[0]
                                    s = prs.slides.add_slide(slide_layout)
                                    s.shapes.title.text = slide["title"]
                                    s.placeholders[1].text = slide["subtitle"]
                            # Chart slides
                            for slide in new_slides:
                                if slide["type"] == "chart":
                                    slide_layout = prs.slide_layouts[5]
                                    s = prs.slides.add_slide(slide_layout)
                                    s.shapes.title.text = slide["caption"]
                                    left = Inches(1)
                                    top = Inches(1.5)
                                    s.shapes.add_picture(BytesIO(slide["image"]), left, top, width=Inches(6))
                            # Summary slide
                            for slide in new_slides:
                                if slide["type"] == "summary":
                                    slide_layout = prs.slide_layouts[1]
                                    s = prs.slides.add_slide(slide_layout)
                                    s.shapes.title.text = slide["title"]
                                    tf = s.placeholders[1].text_frame
                                    for insight in slide["insights"]:
                                        p = tf.add_paragraph()
                                        p.text = insight
                                        p.level = 0
                            pptx_io = BytesIO()
                            prs.save(pptx_io)
                            pptx_io.seek(0)
                            st.success("PowerPoint generated!")

                            # --- Show PowerPoint Preview (Windows/PowerPoint only) ---
                            try:
                                slide_imgs = []
                                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_pptx:
                                    tmp_pptx.write(pptx_io.getvalue())
                                    tmp_pptx.flush()
                                    try:
                                        import comtypes.client
                                        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                                        powerpoint.Visible = 1
                                        presentation = powerpoint.Presentations.Open(tmp_pptx.name)
                                        export_dir = tempfile.mkdtemp()
                                        presentation.SaveAs(export_dir, 17)  # 17 = ppSaveAsPNG
                                        presentation.Close()
                                        powerpoint.Quit()
                                        for fname in sorted(os.listdir(export_dir)):
                                            if fname.endswith(".PNG"):
                                                with open(os.path.join(export_dir, fname), "rb") as imgf:
                                                    slide_imgs.append(imgf.read())
                                        import shutil
                                        shutil.rmtree(export_dir)
                                    except Exception:
                                        st.info("PowerPoint preview is only available on Windows with PowerPoint installed.")
                            except Exception as e:
                                st.info("PowerPoint preview not available: " + str(e))

                            st.download_button(
                                label="Download PowerPoint",
                                data=pptx_io,
                                file_name=f"{file_name}_report.pptx",
                                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                            )
                    else:
                        st.info("Enable 'Report Editing' in the sidebar to customize your report.")

            else:
                st.error("Please upload a valid .xlsx or .csv file.")
                return
        except Exception as e:
            st.error(f"Error reading file: {e}")

# Update process_uploaded_data to accept UI toggles and use expanders for sections
def process_uploaded_data(file_name, file_bytes, is_excel, show_data=True, show_summary=True, show_charts=True, show_qa=True):
    """
    Unified handler for both Excel and CSV uploads.
    Returns: all_images, all_captions, all_insights, summary, insights, images, captions (for CSV)
    """
    all_images = []
    all_captions = []
    all_insights = []
    summary = None
    insights = None
    images = None
    captions = None

    chart_selections = {}
    if is_excel:
        excel_buffer = BytesIO(file_bytes)
        xls = pd.ExcelFile(excel_buffer, engine="openpyxl")
        st.write(f"Detected sheets: {xls.sheet_names}")
        for sheet in xls.sheet_names:
            st.subheader(f"Sheet: {sheet}")
            df = pd.read_excel(xls, sheet_name=sheet, engine="openpyxl")
            # --- Configurable Parameters ---
            num_cols = list(df.select_dtypes(include='number').columns)
            cat_cols = list(df.select_dtypes(include='object').columns)
            with st.expander("Chart Configuration", expanded=False):
                selected_num_cols = st.multiselect(
                    "Numerical columns", num_cols, default=num_cols, key=f"num_cols_{sheet}"
                )
                selected_cat_cols = st.multiselect(
                    "Categorical columns", cat_cols, default=cat_cols, key=f"cat_cols_{sheet}"
                )
                chart_types = st.multiselect(
                    "Chart types", ["hist", "line", "bar", "pie"], default=["hist", "line", "bar", "pie"], key=f"chart_types_{sheet}"
                )
            chart_selections = {
                "num_cols": selected_num_cols,
                "cat_cols": selected_cat_cols,
                "chart_types": chart_types
            }
            # --- Editable DataFrame ---
            if show_data:
                with st.expander("Data Table", expanded=False):
                    edited_df = st.data_editor(df, num_rows="dynamic", key=f"edit_{sheet}")
            else:
                edited_df = df
            summary, insights = data_analysis_module(edited_df)
            if show_summary:
                with st.expander("Summary Statistics", expanded=False):
                    if 'numerical' in summary:
                        st.write(summary['numerical'])
                    if 'categorical' in summary:
                        for col, vc in summary['categorical'].items():
                            st.write(f"Value counts for {col}:")
                            st.write(vc)
            if show_charts:
                with st.expander("Charts", expanded=True):
                    images, captions = plot_and_save(edited_df, sheet, chart_selections)
                    for img, cap in zip(images, captions):
                        st.image(img, caption=cap)
            else:
                images, captions = [], []
            all_insights.extend(insights)
            all_images.extend(images)
            all_captions.extend(captions)
            if show_qa:
                with st.expander("Ask a Question", expanded=False):
                    st.markdown("**Ask a question about this sheet:**")
                    user_query = st.text_input(f"Ask about '{sheet}'", key=f"query_{sheet}")
                    if user_query:
                        try:
                            answer = answer_user_query(edited_df, user_query)
                            if isinstance(answer, pd.DataFrame) or isinstance(answer, pd.Series):
                                st.write(answer)
                            else:
                                st.info(answer)
                        except Exception as e:
                            st.error(f"Error answering question: {e}")
    else:
        csv_buffer = BytesIO(file_bytes)
        df = pd.read_csv(csv_buffer)
        st.subheader(f"CSV File: {file_name}")
        num_cols = list(df.select_dtypes(include='number').columns)
        cat_cols = list(df.select_dtypes(include='object').columns)
        with st.expander("Chart Configuration", expanded=False):
            selected_num_cols = st.multiselect(
                "Numerical columns", num_cols, default=num_cols, key="num_cols_csv"
            )
            selected_cat_cols = st.multiselect(
                "Categorical columns", cat_cols, default=cat_cols, key="cat_cols_csv"
            )
            chart_types = st.multiselect(
                "Chart types", ["hist", "line", "bar", "pie"], default=["hist", "line", "bar", "pie"], key="chart_types_csv"
            )
        chart_selections = {
            "num_cols": selected_num_cols,
            "cat_cols": selected_cat_cols,
            "chart_types": chart_types
        }
        if show_data:
            with st.expander("Data Table", expanded=False):
                edited_df = st.data_editor(df, num_rows="dynamic", key="edit_csv")
        else:
            edited_df = df
        summary, insights = data_analysis_module(edited_df)
        if show_summary:
            with st.expander("Summary Statistics", expanded=False):
                if 'numerical' in summary:
                    st.write(summary['numerical'])
                if 'categorical' in summary:
                    for col, vc in summary['categorical'].items():
                        st.write(f"Value counts for {col}:")
                        st.write(vc)
        if show_charts:
            with st.expander("Charts", expanded=True):
                images, captions = plot_and_save(edited_df, file_name, chart_selections)
                for img, cap in zip(images, captions):
                    st.image(img, caption=cap)
        else:
            images, captions = [], []
        if show_qa:
            with st.expander("Ask a Question", expanded=False):
                st.markdown("**Ask a question about this file:**")
                user_query = st.text_input("Ask about CSV", key="query_csv")
                if user_query:
                    try:
                        answer = answer_user_query(edited_df, user_query)
                        if isinstance(answer, pd.DataFrame) or isinstance(answer, pd.Series):
                            st.write(answer)
                        else:
                            st.info(answer)
                    except Exception as e:
                        st.error(f"Error answering question: {e}")
        all_images = images
        all_captions = captions
        all_insights = insights
    return all_images, all_captions, all_insights, summary, insights, images, captions

if __name__ == "__main__":
    main()
