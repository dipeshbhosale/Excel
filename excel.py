import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from io import BytesIO
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# === Placeholders for advanced modules ===
# (Implement each in its own function for modularity)

def data_upload_module():
    st.header("Step 1: Upload Data")
    uploaded_files = st.file_uploader(
        "Upload Excel (.xlsx) or CSV (.csv) files", type=["xlsx", "csv"], accept_multiple_files=True
    )
    return uploaded_files

def data_cleaning_module(df):
    st.header("Step 2: Data Cleaning Suggestions")
    return df

def data_analysis_module(df):
    st.header("Step 3: Data Analysis & Insights")
    summary = {}
    insights = []
    num_cols = df.select_dtypes(include='number').columns
    cat_cols = df.select_dtypes(include='object').columns

    # Numerical summary
    if len(num_cols) > 0:
        desc = df[num_cols].describe().T
        summary['numerical'] = desc
        for col in num_cols:
            col_data = df[col].dropna()
            if len(col_data) == 0:
                continue
            mean = col_data.mean()
            median = col_data.median()
            std = col_data.std()
            min_ = col_data.min()
            max_ = col_data.max()
            outliers = col_data[(col_data < mean - 3*std) | (col_data > mean + 3*std)]
            if len(outliers) > 0:
                insights.append(f"Column '{col}' has {len(outliers)} outlier(s).")
            insights.append(f"Column '{col}': mean={mean:.2f}, median={median:.2f}, min={min_}, max={max_}.")
    # Categorical summary
    if len(cat_cols) > 0:
        cat_summary = {}
        for col in cat_cols:
            vc = df[col].value_counts()
            cat_summary[col] = vc
            top = vc.index[0] if not vc.empty else None
            if top:
                insights.append(f"Most frequent value in '{col}' is '{top}' ({vc.iloc[0]} times).")
        summary['categorical'] = cat_summary
    return summary, insights

def plot_and_save(df, sheet_name):
    images = []
    captions = []
    num_cols = df.select_dtypes(include='number').columns
    cat_cols = df.select_dtypes(include='object').columns

    # Numerical columns: histograms and line plots
    for col in num_cols:
        fig, ax = plt.subplots()
        sns.histplot(df[col].dropna(), kde=True, ax=ax)
        ax.set_title(f"Distribution of {col}")
        buf = BytesIO()
        plt.savefig(buf, format='png')
        buf.seek(0)
        images.append(buf.read())
        captions.append(f"Distribution of {col} in {sheet_name}.")
        plt.close(fig)

        # Line plot if time-like index
        if pd.api.types.is_datetime64_any_dtype(df.index):
            fig, ax = plt.subplots()
            df[col].plot(ax=ax)
            ax.set_title(f"{col} over time")
            buf = BytesIO()
            plt.savefig(buf, format='png')
            buf.seek(0)
            images.append(buf.read())
            captions.append(f"Trend of {col} over time in {sheet_name}.")
            plt.close(fig)

    # Categorical columns: bar plots and pie charts
    for col in cat_cols:
        vc = df[col].value_counts().head(10)
        if len(vc) > 1:
            fig, ax = plt.subplots()
            sns.barplot(x=vc.values, y=vc.index, ax=ax)
            ax.set_title(f"Top categories in {col}")
            buf = BytesIO()
            plt.savefig(buf, format='png')
            buf.seek(0)
            images.append(buf.read())
            captions.append(f"Top categories in {col} in {sheet_name}.")
            plt.close(fig)

            fig, ax = plt.subplots()
            vc.plot.pie(autopct='%1.1f%%', ax=ax)
            ax.set_ylabel('')
            ax.set_title(f"Category distribution in {col}")
            buf = BytesIO()
            plt.savefig(buf, format='png')
            buf.seek(0)
            images.append(buf.read())
            captions.append(f"Category distribution in {col} in {sheet_name}.")
            plt.close(fig)
    return images, captions

def add_title_slide(prs, file_name):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Data Analysis Report"
    subtitle.text = f"File: {file_name}\nDate: {datetime.now().strftime('%Y-%m-%d')}"

def add_chart_slide(prs, image_bytes, caption):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = caption
    left = Inches(1)
    top = Inches(1.5)
    pic = slide.shapes.add_picture(BytesIO(image_bytes), left, top, width=Inches(6))
    # Optionally add caption as text box
    # ...

def add_summary_slide(prs, insights):
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Findings"
    tf = slide.placeholders[1].text_frame
    for insight in insights:
        p = tf.add_paragraph()
        p.text = insight
        p.level = 0

def create_pptx(file_name, all_images, all_captions, insights):
    prs = Presentation()
    add_title_slide(prs, file_name)
    for img, cap in zip(all_images, all_captions):
        add_chart_slide(prs, img, cap)
    add_summary_slide(prs, insights)
    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

def reporting_module(selected_charts, insights):
    st.header("Step 5: Reporting & Export")
    st.info("Reporting module is under development.")
    # Placeholder for reporting functionality

def collaboration_module():
    st.header("Step 6: Collaboration & Sharing")
    st.info("Collaboration module is under development.")
    # Placeholder for collaboration functionality

def ai_ml_module(df):
    st.header("Step 7: AI/ML Features")
    st.info("AI/ML module is under development.")
    # Placeholder for AI/ML functionality

def main():
    st.set_page_config(page_title="Data-Driven SaaS Assistant", layout="wide")
    st.title("Excel Data-Driven SaaS Assistant")
    uploaded_file = st.file_uploader("Upload an Excel (.xlsx) or CSV (.csv) file", type=["xlsx", "csv"])
    if uploaded_file is not None:
        try:
            file_name = uploaded_file.name
            file_bytes = uploaded_file.read()
            if not file_bytes:
                st.error("Uploaded file is empty.")
                return
            if file_name.lower().endswith(".xlsx"):
                excel_buffer = BytesIO(file_bytes)
                xls = pd.ExcelFile(excel_buffer, engine="openpyxl")
                st.write(f"Detected sheets: {xls.sheet_names}")
                all_images = []
                all_captions = []
                all_insights = []
                for sheet in xls.sheet_names:
                    st.subheader(f"Sheet: {sheet}")
                    df = pd.read_excel(xls, sheet_name=sheet, engine="openpyxl")
                    st.dataframe(df.head())
                    summary, insights = data_analysis_module(df)
                    all_insights.extend(insights)
                    st.write("Summary statistics:")
                    if 'numerical' in summary:
                        st.write(summary['numerical'])
                    if 'categorical' in summary:
                        for col, vc in summary['categorical'].items():
                            st.write(f"Value counts for {col}:")
                            st.write(vc)
                    images, captions = plot_and_save(df, sheet)
                    for img, cap in zip(images, captions):
                        st.image(img, caption=cap)
                    all_images.extend(images)
                    all_captions.extend(captions)
                if st.button("Generate PowerPoint Report"):
                    pptx_io = create_pptx(file_name, all_images, all_captions, all_insights)
                    st.success("PowerPoint generated!")
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_io,
                        file_name=f"{file_name}_report.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            elif file_name.lower().endswith(".csv"):
                csv_buffer = BytesIO(file_bytes)
                df = pd.read_csv(csv_buffer)
                st.subheader(f"CSV File: {file_name}")
                st.dataframe(df.head())
                summary, insights = data_analysis_module(df)
                st.write("Summary statistics:")
                if 'numerical' in summary:
                    st.write(summary['numerical'])
                if 'categorical' in summary:
                    for col, vc in summary['categorical'].items():
                        st.write(f"Value counts for {col}:")
                        st.write(vc)
                images, captions = plot_and_save(df, file_name)
                for img, cap in zip(images, captions):
                    st.image(img, caption=cap)
                if st.button("Generate PowerPoint Report"):
                    pptx_io = create_pptx(file_name, images, captions, insights)
                    st.success("PowerPoint generated!")
                    st.download_button(
                        label="Download PowerPoint",
                        data=pptx_io,
                        file_name=f"{file_name}_report.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                    )
            else:
                st.error("Please upload a valid .xlsx or .csv file.")
                return
        except Exception as e:
            st.error(f"Error reading file: {e}")

if __name__ == "__main__":
    main()